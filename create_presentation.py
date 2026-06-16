#!/usr/bin/env python3
"""Create the final Company A telecom churn business proposal deck."""

from __future__ import annotations

import json
from pathlib import Path

import joblib
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.metrics import auc, confusion_matrix, roc_curve
from sklearn.model_selection import train_test_split

ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "presentation_assets"
OUTPUT_PATH = ROOT / "Company_A_Churn_Proposal.pptx"

NAVY = "102A43"
BLUE = "1769AA"
CYAN = "2CB1BC"
TEAL = "00A896"
GREEN = "2D9D78"
ORANGE = "F59E0B"
RED = "D64545"
INK = "243B53"
GRAY = "627D98"
LIGHT = "F4F7FA"
PALE_BLUE = "E8F1F8"
WHITE = "FFFFFF"
FONT = "Liberation Sans"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def load_analysis() -> dict:
    client = pd.read_csv(ROOT / "telecom" / "Client.csv")
    record = pd.read_csv(ROOT / "telecom" / "Record.csv")
    data = record.merge(client, on="Customer_ID", validate="one_to_one")
    features = data.drop(columns=["Customer_ID", "churn"])
    target = data["churn"].astype(int)

    _, temp_features, _, temp_target = train_test_split(
        features,
        target,
        test_size=0.4,
        random_state=42,
        stratify=target,
    )
    _, test_features, _, test_target = train_test_split(
        temp_features,
        temp_target,
        test_size=0.5,
        random_state=42,
        stratify=temp_target,
    )
    model = joblib.load(ROOT / "artifacts" / "advanced" / "advanced_churn_model.joblib")
    probability = model.predict_proba(test_features)[:, 1]

    scored = test_features.copy()
    scored["churn"] = test_target.to_numpy()
    scored["probability"] = probability
    scored = scored.sort_values("probability", ascending=False).reset_index(drop=True)

    top_twenty = scored.head(int(len(scored) * 0.2))
    capture = top_twenty["churn"].sum() / scored["churn"].sum()
    top_rate = top_twenty["churn"].mean()
    overall_rate = scored["churn"].mean()
    average_target_revenue = top_twenty["rev_Mean"].mean()

    annual_customers = 100_000
    targeted_customers = int(annual_customers * 0.2)
    expected_churners = target.mean() * annual_customers * capture
    action_cost = 15
    save_rate = 0.20
    gross_value = expected_churners * save_rate * average_target_revenue * 12
    campaign_cost = targeted_customers * action_cost
    net_value = gross_value - campaign_cost

    return {
        "data": data,
        "features": features,
        "target": target,
        "test_features": test_features,
        "test_target": test_target,
        "probability": probability,
        "scored": scored,
        "capture": capture,
        "top_rate": top_rate,
        "overall_rate": overall_rate,
        "lift": top_rate / overall_rate,
        "average_target_revenue": average_target_revenue,
        "gross_value": gross_value,
        "campaign_cost": campaign_cost,
        "net_value": net_value,
        "roi": net_value / campaign_cost,
        "metrics": json.loads(
            (ROOT / "artifacts" / "advanced" / "metrics.json").read_text()
        ),
    }


def chart_style() -> None:
    sns.set_theme(style="whitegrid")
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "axes.titlesize": 14,
            "axes.labelsize": 11,
            "xtick.labelsize": 9,
            "ytick.labelsize": 9,
            "axes.titleweight": "bold",
        }
    )


def save_eda_chart(analysis: dict) -> None:
    data = analysis["data"]
    fig, axes = plt.subplots(1, 3, figsize=(12, 3.6))
    specifications = [
        ("eqpdays", "Equipment age", "days"),
        ("change_mou", "Usage momentum", "change in minutes"),
        ("mou_Mean", "Monthly usage", "minutes"),
    ]
    for axis, (column, title, unit) in zip(axes, specifications, strict=True):
        frame = data[[column, "churn"]].dropna().copy()
        frame["quartile"] = pd.qcut(
            frame[column], 4, labels=["Q1", "Q2", "Q3", "Q4"], duplicates="drop"
        )
        rates = frame.groupby("quartile", observed=True)["churn"].mean() * 100
        colors = [rgb_value / 255 for rgb_value in rgb(CYAN)]
        axis.bar(rates.index.astype(str), rates.values, color=f"#{CYAN}")
        axis.set_ylim(0, max(65, rates.max() + 7))
        axis.set_title(title)
        axis.set_ylabel("Churn rate (%)")
        axis.set_xlabel(f"Quartiles by {unit}")
        axis.grid(axis="x", visible=False)
        for position, value in enumerate(rates.values):
            axis.text(position, value + 1, f"{value:.1f}%", ha="center", fontsize=9)
    fig.suptitle("Churn changes materially across actionable customer signals", y=1.04)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "eda_signals.png", dpi=220, bbox_inches="tight")
    plt.close(fig)


def save_model_chart() -> None:
    baseline = pd.read_csv(ROOT / "artifacts" / "validation_metrics.csv")
    advanced = pd.read_csv(ROOT / "artifacts" / "advanced_tuning" / "results.csv")
    values = pd.DataFrame(
        {
            "Model": [
                "Logistic regression",
                "Random forest",
                "Histogram boosting",
                "XGBoost + features",
                "LightGBM + features",
                "Final ensemble",
            ],
            "ROC-AUC": [
                baseline.loc[
                    baseline["model"] == "logistic_regression", "roc_auc"
                ].iloc[0],
                baseline.loc[baseline["model"] == "random_forest", "roc_auc"].iloc[0],
                baseline.loc[
                    baseline["model"] == "hist_gradient_boosting", "roc_auc"
                ].iloc[0],
                advanced.loc[advanced["name"] == "xgb_d5", "roc_auc"].iloc[0],
                advanced.loc[advanced["name"] == "lgbm_63", "roc_auc"].iloc[0],
                0.7079970583,
            ],
        }
    ).sort_values("ROC-AUC")
    colors = [f"#{GRAY}"] * (len(values) - 1) + [f"#{TEAL}"]
    fig, axis = plt.subplots(figsize=(8.6, 4.3))
    bars = axis.barh(values["Model"], values["ROC-AUC"], color=colors)
    axis.set_xlim(0.58, 0.72)
    axis.set_xlabel("Validation ROC-AUC")
    axis.set_title("Feature engineering and ensembling produced the strongest ranking model")
    axis.grid(axis="y", visible=False)
    for bar, value in zip(bars, values["ROC-AUC"], strict=True):
        axis.text(value + 0.002, bar.get_y() + bar.get_height() / 2, f"{value:.3f}", va="center")
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "model_benchmark.png", dpi=220, bbox_inches="tight")
    plt.close(fig)


def save_validation_chart(analysis: dict) -> None:
    target = analysis["test_target"]
    probability = analysis["probability"]
    fpr, tpr, _ = roc_curve(target, probability)
    score = auc(fpr, tpr)
    matrix = confusion_matrix(target, probability >= 0.5)

    fig, axes = plt.subplots(1, 2, figsize=(9.5, 4))
    axes[0].plot(fpr, tpr, color=f"#{BLUE}", linewidth=3)
    axes[0].plot([0, 1], [0, 1], linestyle="--", color=f"#{GRAY}")
    axes[0].fill_between(fpr, tpr, alpha=0.12, color=f"#{BLUE}")
    axes[0].set_title(f"Held-out ROC curve | AUC = {score:.3f}")
    axes[0].set_xlabel("False positive rate")
    axes[0].set_ylabel("True positive rate")
    sns.heatmap(
        matrix,
        annot=True,
        fmt=",d",
        cmap=sns.light_palette(f"#{TEAL}", as_cmap=True),
        cbar=False,
        ax=axes[1],
        xticklabels=["Stay", "Churn"],
        yticklabels=["Stay", "Churn"],
    )
    axes[1].set_title("Confusion matrix | threshold 0.50")
    axes[1].set_xlabel("Predicted")
    axes[1].set_ylabel("Actual")
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "model_validation.png", dpi=220, bbox_inches="tight")
    plt.close(fig)


def friendly_feature_name(raw_name: str) -> str:
    name = raw_name.replace("numeric__", "").replace("categorical__", "")
    mapping = {
        "eqpdays": "Equipment age",
        "months": "Customer tenure",
        "current_to_3m_usage": "Current vs 3-month usage",
        "current_to_lifetime_usage": "Current vs lifetime usage",
        "recurring_revenue_share": "Recurring revenue share",
        "equipment_age_to_tenure": "Equipment age / tenure",
        "current_to_3m_revenue": "Current vs 3-month revenue",
        "hnd_price": "Handset price",
        "totmrc_Mean": "Monthly recurring charge",
        "crclscod": "Credit class",
        "voice_drop_rate": "Voice drop rate",
        "minutes_per_call": "Minutes per call",
    }
    return mapping.get(name, name.replace("_", " ").title())


def save_feature_chart() -> None:
    importance = pd.read_csv(ROOT / "artifacts" / "advanced" / "feature_importance.csv").head(12)
    importance = importance.iloc[::-1].copy()
    importance["label"] = importance["feature"].map(friendly_feature_name)
    importance["share"] = importance["gain"] / importance["gain"].sum() * 100
    fig, axis = plt.subplots(figsize=(8.4, 4.7))
    colors = [f"#{CYAN}"] * 6 + [f"#{BLUE}"] * 6
    axis.barh(importance["label"], importance["share"], color=colors)
    axis.set_xlabel("Share of gain among top 12 features (%)")
    axis.set_title("The model combines lifecycle, engagement, value and service-quality signals")
    axis.grid(axis="y", visible=False)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "feature_importance.png", dpi=220, bbox_inches="tight")
    plt.close(fig)


def save_targeting_chart(analysis: dict) -> None:
    scored = analysis["scored"]
    shares = np.arange(0.05, 1.01, 0.05)
    captures = []
    lifts = []
    for share in shares:
        top = scored.head(int(len(scored) * share))
        captures.append(top["churn"].sum() / scored["churn"].sum())
        lifts.append(top["churn"].mean() / scored["churn"].mean())

    fig, axes = plt.subplots(1, 2, figsize=(9.5, 4))
    axes[0].plot(shares * 100, np.array(captures) * 100, color=f"#{TEAL}", linewidth=3)
    axes[0].plot([0, 100], [0, 100], linestyle="--", color=f"#{GRAY}")
    axes[0].scatter([20], [analysis["capture"] * 100], color=f"#{ORANGE}", s=80, zorder=3)
    axes[0].annotate(
        f"Top 20% captures\n{analysis['capture']:.0%} of churners",
        (20, analysis["capture"] * 100),
        xytext=(32, 22),
        arrowprops={"arrowstyle": "->", "color": f"#{INK}"},
    )
    axes[0].set_title("Cumulative churn capture")
    axes[0].set_xlabel("Customers targeted (%)")
    axes[0].set_ylabel("Churners captured (%)")

    axes[1].plot(shares * 100, lifts, color=f"#{BLUE}", linewidth=3)
    axes[1].axhline(1, linestyle="--", color=f"#{GRAY}")
    axes[1].scatter([20], [analysis["lift"]], color=f"#{ORANGE}", s=80, zorder=3)
    axes[1].annotate(
        f"{analysis['lift']:.2f}x lift",
        (20, analysis["lift"]),
        xytext=(35, 1.58),
        arrowprops={"arrowstyle": "->", "color": f"#{INK}"},
    )
    axes[1].set_title("Risk concentration")
    axes[1].set_xlabel("Customers targeted (%)")
    axes[1].set_ylabel("Churn rate lift vs average")
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "targeting_performance.png", dpi=220, bbox_inches="tight")
    plt.close(fig)


def save_impact_chart(analysis: dict) -> None:
    expected_churners = (
        analysis["overall_rate"] * 100_000 * analysis["capture"]
    )
    values = []
    for save_rate in (0.10, 0.20, 0.30):
        gross = expected_churners * save_rate * analysis["average_target_revenue"] * 12
        values.append((gross - analysis["campaign_cost"]) / 1_000_000)
    fig, axis = plt.subplots(figsize=(7.5, 4.1))
    bars = axis.bar(
        ["Conservative\n10% saved", "Base\n20% saved", "Upside\n30% saved"],
        values,
        color=[f"#{PALE_BLUE}", f"#{TEAL}", f"#{BLUE}"],
        edgecolor=f"#{BLUE}",
    )
    axis.axhline(0, color=f"#{INK}", linewidth=1)
    axis.set_ylabel("Illustrative net annual value ($M)")
    axis.set_title("Retention value remains positive across intervention scenarios")
    axis.grid(axis="x", visible=False)
    for bar, value in zip(bars, values, strict=True):
        axis.text(
            bar.get_x() + bar.get_width() / 2,
            value + 0.06,
            f"${value:.2f}M",
            ha="center",
            fontweight="bold",
        )
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "impact_scenarios.png", dpi=220, bbox_inches="tight")
    plt.close(fig)


def generate_assets(analysis: dict) -> None:
    ASSET_DIR.mkdir(exist_ok=True)
    chart_style()
    save_eda_chart(analysis)
    save_model_chart()
    save_validation_chart(analysis)
    save_feature_chart()
    save_targeting_chart(analysis)
    save_impact_chart(analysis)


def set_background(slide, color: str = WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, pieces, x, y, w, h, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.04)
    paragraph = frame.paragraphs[0]
    for text, bold, piece_color in pieces:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(piece_color or color)
    return box


def add_rect(slide, x, y, w, h, fill, radius=True, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_line(slide, x1, y1, x2, y2, color=GRAY, width=1.5):
    line = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.65, 0.35, 11.9, 0.55, 25, NAVY, True)
    add_rect(slide, 0.65, 0.98, 0.72, 0.06, TEAL, radius=False)
    if subtitle:
        add_text(slide, subtitle, 1.52, 0.88, 10.9, 0.3, 10.5, GRAY)


def add_footer(slide, number: int, source: str | None = None) -> None:
    add_line(slide, 0.65, 7.12, 12.68, 7.12, PALE_BLUE, 1)
    if source:
        add_text(slide, source, 0.68, 7.18, 10.9, 0.18, 7.5, GRAY)
    add_text(slide, f"{number:02d}", 12.0, 7.15, 0.55, 0.2, 8, GRAY, True, PP_ALIGN.RIGHT)


def add_metric_card(slide, value, label, x, y, w=2.2, color=TEAL, note=None):
    add_rect(slide, x, y, w, 1.25, WHITE, line=PALE_BLUE)
    add_rect(slide, x, y, 0.08, 1.25, color, radius=False)
    add_text(slide, value, x + 0.22, y + 0.2, w - 0.35, 0.48, 25, color, True)
    add_text(slide, label, x + 0.22, y + 0.72, w - 0.35, 0.3, 10.5, INK, True)
    if note:
        add_text(slide, note, x + 0.22, y + 1.0, w - 0.35, 0.17, 7.5, GRAY)


def add_bullet_list(slide, items, x, y, w, h, size=16, color=INK, accent=TEAL):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.03)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = rgb(color)
        paragraph.space_after = Pt(9)
        paragraph.level = 0
        paragraph.text = f"•  {item}"
    return box


def add_picture(slide, path: Path, x, y, w, h=None):
    kwargs = {"width": Inches(w)}
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def add_process_step(slide, number, title, body, x, y, w, color=BLUE):
    add_rect(slide, x, y, w, 1.5, WHITE, line=PALE_BLUE)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.25), Inches(0.48), Inches(0.48)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(color)
    circle.line.color.rgb = rgb(color)
    add_text(slide, str(number), x + 0.18, y + 0.29, 0.48, 0.28, 12, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.82, y + 0.2, w - 1.0, 0.35, 14, NAVY, True)
    add_text(slide, body, x + 0.82, y + 0.62, w - 1.02, 0.65, 10.5, GRAY)


def build_deck(analysis: dict) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    # 1. Cover
    slide = presentation.slides.add_slide(blank)
    set_background(slide, NAVY)
    add_rect(slide, 8.95, 0, 4.38, 7.5, BLUE, radius=False)
    for x, y, radius, color in [
        (9.35, 0.8, 1.5, CYAN),
        (11.3, 2.0, 0.85, ORANGE),
        (9.8, 4.2, 2.2, TEAL),
        (11.9, 5.7, 0.7, WHITE),
    ]:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(radius), Inches(radius)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(color)
        circle.fill.transparency = 18 if color != WHITE else 65
        circle.line.fill.background()
    add_text(slide, "COMPANY A | CUSTOMER RETENTION PoC", 0.75, 0.72, 7.8, 0.35, 11, CYAN, True)
    add_text(slide, "Retain the Right\nCustomers Before\nThey Leave", 0.75, 1.42, 7.7, 2.2, 34, WHITE, True)
    add_text(
        slide,
        "An AI-powered churn prevention proposal built from 100,000 customer records",
        0.78,
        4.08,
        7.3,
        0.72,
        17,
        "D9EAF4",
    )
    add_rect(slide, 0.78, 5.18, 6.85, 1.05, "173F5F")
    add_rich_text(
        slide,
        [
            ("Recommendation: ", True, CYAN),
            ("rank customers weekly and focus retention actions on the highest-risk 20%.", False, WHITE),
        ],
        1.0,
        5.46,
        6.4,
        0.55,
        15,
    )
    add_text(slide, "GCI World 2026 Final Assignment", 0.78, 6.86, 4.2, 0.24, 9, "AFC7D5")

    # 2. Executive recommendation
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Executive recommendation", "Move from broad retention campaigns to risk-ranked, action-specific outreach")
    add_metric_card(slide, "0.702", "Held-out ROC-AUC", 0.75, 1.45, 2.45, BLUE, "20,000-customer test set")
    add_metric_card(slide, "30.0%", "Churners captured", 3.45, 1.45, 2.45, TEAL, "by targeting the top 20%")
    add_metric_card(slide, "1.50×", "Risk lift", 6.15, 1.45, 2.45, ORANGE, "74.4% vs 49.6% baseline")
    add_metric_card(slide, "$1.65M", "Illustrative net value", 8.85, 1.45, 2.75, GREEN, "per 100k customers/year")
    add_rect(slide, 0.75, 3.15, 11.8, 2.75, LIGHT, line=PALE_BLUE)
    add_text(slide, "What Company A should do", 1.05, 3.48, 3.3, 0.4, 19, NAVY, True)
    add_bullet_list(
        slide,
        [
            "Score active customers weekly using usage, billing, handset and service-quality signals.",
            "Prioritize the highest-risk 20% for targeted offers, service recovery or handset upgrades.",
            "Measure incremental retention through a randomized control group before full rollout.",
        ],
        1.0,
        4.02,
        5.35,
        1.5,
        14,
    )
    add_text(slide, "Why now", 6.75, 3.48, 2.2, 0.4, 19, NAVY, True)
    add_bullet_list(
        slide,
        [
            "The data already contains predictive behavior before churn.",
            "The solution is deployable as a probability-ranked workflow, not a one-time model.",
            "The business case remains positive under conservative intervention assumptions.",
        ],
        6.7,
        4.02,
        5.15,
        1.5,
        14,
    )
    add_footer(slide, 2, "Source: Company A dataset and held-out model evaluation; financial impact is an illustrative scenario.")

    # 3. Market context
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Retention is a telecom growth problem", "Saturated markets increase the value of protecting the installed customer base")
    add_rect(slide, 0.75, 1.35, 5.55, 4.95, NAVY)
    add_text(slide, "The strategic tension", 1.05, 1.7, 3.5, 0.45, 20, CYAN, True)
    add_text(
        slide,
        "Telecom operators must defend revenue while customers face low switching friction, aggressive offers and rising service expectations.",
        1.05,
        2.3,
        4.9,
        1.1,
        20,
        WHITE,
        True,
    )
    add_text(slide, "Industry signal", 1.05, 3.83, 2.4, 0.35, 12, "AFC7D5", True)
    add_rich_text(
        slide,
        [
            ("22% ", True, ORANGE),
            ("average annual churn across providers in Deloitte's 2026 outlook.", False, WHITE),
        ],
        1.05,
        4.25,
        4.75,
        0.78,
        17,
    )
    add_text(
        slide,
        "Company A's response should be selective: identify preventable churn early and spend only where intervention is economically justified.",
        1.05,
        5.38,
        4.85,
        0.65,
        13,
        "D9EAF4",
    )
    cards = [
        ("Revenue risk", "Customers leave with recurring revenue and future lifetime value.", RED),
        ("Service signals", "Usage decline, failures and support activity can reveal friction.", ORANGE),
        ("Action window", "The target records churn 31–60 days after observation.", TEAL),
    ]
    for index, (heading, body, color) in enumerate(cards):
        y = 1.48 + index * 1.58
        add_rect(slide, 6.7, y, 5.55, 1.28, WHITE, line=PALE_BLUE)
        add_rect(slide, 6.7, y, 0.1, 1.28, color, radius=False)
        add_text(slide, heading, 7.05, y + 0.18, 2.0, 0.3, 14, NAVY, True)
        add_text(slide, body, 7.05, y + 0.56, 4.75, 0.48, 11.5, GRAY)
    add_rect(slide, 6.7, 6.07, 5.55, 0.46, PALE_BLUE, line=PALE_BLUE)
    add_text(slide, "Business question: who should receive which retention action, and when?", 6.95, 6.17, 5.0, 0.22, 11, BLUE, True)
    add_footer(slide, 3, "Sources: [1] Deloitte 2026 Telecommunications Industry Outlook; [2] Company A dataset description.")

    # 4. Data overview
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Company A already has the data required for a retention PoC")
    add_metric_card(slide, "100,000", "Customers", 0.75, 1.35, 2.4, BLUE)
    add_metric_card(slide, "98", "Raw predictors", 3.42, 1.35, 2.4, CYAN)
    add_metric_card(slide, "2", "Linked tables", 6.09, 1.35, 2.4, ORANGE)
    add_metric_card(slide, "31–60d", "Prediction window", 8.76, 1.35, 2.65, TEAL)
    categories = [
        ("Usage", "Minutes, calls, peak/off-peak behavior, recent change"),
        ("Billing", "Revenue, recurring charges, overage and roaming"),
        ("Experience", "Dropped, blocked and unanswered calls; care contacts"),
        ("Lifecycle", "Tenure, equipment age, handset and household accounts"),
        ("Customer", "Area, credit class and available demographic attributes"),
    ]
    for index, (heading, body) in enumerate(categories):
        x = 0.8 + (index % 3) * 4.0
        y = 3.2 + (index // 3) * 1.42
        width = 3.55 if index < 3 else 5.55
        if index >= 3:
            x = 0.8 + (index - 3) * 5.95
        add_rect(slide, x, y, width, 1.1, LIGHT, line=PALE_BLUE)
        add_text(slide, heading, x + 0.2, y + 0.17, 1.2, 0.3, 13, BLUE, True)
        add_text(slide, body, x + 1.3, y + 0.14, width - 1.5, 0.67, 10.5, GRAY)
    add_rect(slide, 0.8, 6.12, 11.45, 0.5, "FFF4E5", line="F7D9A6")
    add_text(
        slide,
        "Important caveat: the supplied target is almost 50/50 churn, unlike many live telecom portfolios. Production calibration must use Company A's current base rate.",
        1.0,
        6.24,
        11.05,
        0.25,
        10.5,
        "8A5A00",
        True,
    )
    add_footer(slide, 4, "Source: Client.csv and Record.csv; one-to-one merge on Customer_ID.")

    # 5. EDA
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "EDA reveals actionable warning signals", "Single variables are insufficient, but several behaviors show consistent directional risk")
    add_picture(slide, ASSET_DIR / "eda_signals.png", 0.65, 1.35, 8.45)
    add_rect(slide, 9.2, 1.55, 3.25, 4.85, NAVY)
    add_text(slide, "What the patterns say", 9.55, 1.9, 2.45, 0.4, 18, CYAN, True)
    add_bullet_list(
        slide,
        [
            "Older equipment is associated with higher churn.",
            "Declining usage is an early warning signal.",
            "Low current usage indicates weaker engagement.",
            "No single variable separates churners; multivariate modeling is required.",
        ],
        9.48,
        2.58,
        2.45,
        2.55,
        11.8,
        WHITE,
    )
    add_text(slide, "Actionable hypothesis", 9.55, 5.43, 2.1, 0.28, 10.5, "AFC7D5", True)
    add_text(
        slide,
        "Combine lifecycle, engagement and service signals to trigger targeted actions.",
        9.55,
        5.76,
        2.35,
        0.5,
        10.5,
        WHITE,
        True,
    )
    add_footer(slide, 5, "Source: Company A EDA; quartiles calculated after excluding missing values for each feature.")

    # 6. Problem statement
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Problem statement and success definition")
    add_rect(slide, 0.75, 1.35, 11.8, 1.32, NAVY)
    add_text(slide, "Business problem", 1.05, 1.65, 2.1, 0.35, 13, CYAN, True)
    add_text(
        slide,
        "Company A cannot economically treat every customer, but it lacks a repeatable method to identify who is most likely to churn and which operational lever should be used.",
        3.05,
        1.57,
        8.75,
        0.65,
        18,
        WHITE,
        True,
    )
    columns = [
        ("ML objective", "Rank customers by probability of churn 31–60 days ahead.", BLUE),
        ("Business objective", "Concentrate intervention spend where avoidable revenue risk is highest.", TEAL),
        ("Success metrics", "ROC-AUC and PR-AUC for ranking; capture, lift and incremental retention for business value.", ORANGE),
    ]
    for index, (heading, body, color) in enumerate(columns):
        x = 0.78 + index * 4.02
        add_rect(slide, x, 3.12, 3.7, 2.3, LIGHT, line=PALE_BLUE)
        add_rect(slide, x, 3.12, 3.7, 0.14, color, radius=False)
        add_text(slide, heading, x + 0.24, 3.48, 3.1, 0.38, 16, NAVY, True)
        add_text(slide, body, x + 0.24, 4.05, 3.15, 0.92, 13, GRAY)
    add_rect(slide, 1.6, 5.85, 10.0, 0.65, PALE_BLUE, line=PALE_BLUE)
    add_text(
        slide,
        "Decision rule for the PoC: use risk ranking to define campaign capacity; do not treat the model threshold as an automatic cancellation or denial decision.",
        1.9,
        6.02,
        9.45,
        0.3,
        12,
        BLUE,
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 6, "Source: Project problem formulation based on Company A's churn label and available intervention signals.")

    # 7. Solution architecture
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Solution: a weekly risk-ranking and action engine")
    steps = [
        ("Ingest", "Refresh customer, usage and billing records."),
        ("Engineer", "Create 34 lifecycle, trend, value and service ratios."),
        ("Score", "Blend LightGBM with native-categorical boosting."),
        ("Prioritize", "Rank customers and apply campaign capacity."),
        ("Act & learn", "Deliver offers; measure uplift against controls."),
    ]
    for index, (heading, body) in enumerate(steps):
        x = 0.55 + index * 2.55
        add_process_step(slide, index + 1, heading, body, x, 1.75, 2.25, [BLUE, CYAN, TEAL, ORANGE, GREEN][index])
        if index < 4:
            add_text(slide, "→", x + 2.27, 2.18, 0.35, 0.4, 20, GRAY, True, PP_ALIGN.CENTER)
    add_rect(slide, 0.75, 3.85, 11.8, 2.25, LIGHT, line=PALE_BLUE)
    add_text(slide, "Why this design is practical", 1.05, 4.18, 3.2, 0.4, 18, NAVY, True)
    add_bullet_list(
        slide,
        [
            "Works on existing tabular data; no new platform is required for the PoC.",
            "Produces probabilities that support multiple campaign capacities and cost structures.",
        ],
        1.0,
        4.72,
        5.15,
        1.05,
        13,
    )
    add_text(slide, "Controls built into the workflow", 6.55, 4.18, 3.4, 0.4, 18, NAVY, True)
    add_bullet_list(
        slide,
        [
            "Human-approved action catalogue; no adverse automated decisions.",
            "Monitoring for drift, calibration, fairness and incremental campaign lift.",
        ],
        6.5,
        4.72,
        5.3,
        1.05,
        13,
    )
    add_footer(slide, 7, "Source: Implemented pipeline in telecom_features.py and train_advanced_model.py.")

    # 8. Model benchmark
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Model development moved beyond the tutorial baseline")
    add_picture(slide, ASSET_DIR / "model_benchmark.png", 0.65, 1.35, 8.55)
    add_rect(slide, 9.25, 1.55, 3.1, 4.85, LIGHT, line=PALE_BLUE)
    add_text(slide, "Final model", 9.55, 1.9, 1.8, 0.35, 17, NAVY, True)
    add_text(slide, "65%", 9.55, 2.48, 0.95, 0.5, 25, BLUE, True)
    add_text(slide, "Engineered LightGBM", 10.45, 2.56, 1.4, 0.4, 11.5, INK, True)
    add_text(slide, "35%", 9.55, 3.22, 0.95, 0.5, 25, TEAL, True)
    add_text(slide, "Categorical ensemble", 10.45, 3.30, 1.4, 0.4, 11.5, INK, True)
    add_line(slide, 9.55, 3.98, 11.95, 3.98, PALE_BLUE, 1)
    add_text(slide, "Selection protocol", 9.55, 4.25, 2.0, 0.3, 12, GRAY, True)
    add_bullet_list(
        slide,
        [
            "60% train",
            "20% validation",
            "20% untouched test",
            "Selection by validation ROC-AUC",
        ],
        9.52,
        4.65,
        2.1,
        1.35,
        11.5,
    )
    add_footer(slide, 8, "Source: Reproducible model runs; test data excluded from feature, model and blend selection.")

    # 9. Validation
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "The final model generalizes to unseen customers")
    add_picture(slide, ASSET_DIR / "model_validation.png", 0.65, 1.35, 8.65)
    metrics = analysis["metrics"]["test_metrics_at_0_5"]
    add_metric_card(slide, f"{metrics['roc_auc']:.3f}", "ROC-AUC", 9.55, 1.52, 2.45, BLUE)
    add_metric_card(slide, f"{metrics['pr_auc']:.3f}", "PR-AUC", 9.55, 2.98, 2.45, TEAL)
    add_metric_card(slide, f"{metrics['accuracy']:.1%}", "Accuracy", 9.55, 4.44, 2.45, ORANGE)
    add_rect(slide, 9.45, 5.92, 2.8, 0.55, "FFF4E5", line="F7D9A6")
    add_text(slide, "Use ranking, not accuracy alone", 9.63, 6.08, 2.42, 0.22, 9.5, "8A5A00", True, PP_ALIGN.CENTER)
    add_footer(slide, 9, "Source: Held-out test set, n=20,000. Confusion matrix shown at a 0.50 probability threshold.")

    # 10. Drivers
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Churn risk is driven by interpretable business signals")
    add_picture(slide, ASSET_DIR / "feature_importance.png", 0.65, 1.35, 8.45)
    add_rect(slide, 9.15, 1.48, 3.25, 5.15, NAVY)
    add_text(slide, "Four intervention themes", 9.48, 1.82, 2.5, 0.38, 17, CYAN, True)
    themes = [
        ("Lifecycle", "Equipment age and tenure", BLUE),
        ("Engagement", "Current vs historical usage", CYAN),
        ("Value", "Recurring revenue and handset value", ORANGE),
        ("Experience", "Voice drops and call completion", RED),
    ]
    for index, (heading, body, color) in enumerate(themes):
        y = 2.55 + index * 0.91
        add_rect(slide, 9.45, y, 0.13, 0.65, color, radius=False)
        add_text(slide, heading, 9.75, y, 1.15, 0.25, 11.5, WHITE, True)
        add_text(slide, body, 10.85, y, 1.2, 0.48, 9.5, "D9EAF4")
    add_text(slide, "Interpretation guardrail", 9.48, 6.0, 2.2, 0.24, 10, "AFC7D5", True)
    add_text(slide, "Importance identifies prediction signals, not causal effects.", 9.48, 6.3, 2.35, 0.28, 9.5, WHITE)
    add_footer(slide, 10, "Source: LightGBM gain importance. Causal impact must be validated through controlled interventions.")

    # 11. Targeting strategy
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "A top-20% strategy concentrates campaign spend")
    add_picture(slide, ASSET_DIR / "targeting_performance.png", 0.65, 1.35, 8.8)
    add_rect(slide, 9.6, 1.52, 2.55, 4.85, LIGHT, line=PALE_BLUE)
    add_text(slide, "Recommended PoC capacity", 9.88, 1.88, 1.95, 0.65, 15, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, "20%", 9.9, 2.75, 1.95, 0.65, 31, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, "of customers", 9.9, 3.38, 1.95, 0.3, 12, GRAY, True, PP_ALIGN.CENTER)
    add_line(slide, 9.93, 3.92, 11.82, 3.92, PALE_BLUE, 1)
    add_text(slide, "74.4%", 9.9, 4.22, 1.95, 0.5, 24, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "observed churn rate", 9.9, 4.75, 1.95, 0.35, 10.5, GRAY, True, PP_ALIGN.CENTER)
    add_text(slide, "1.50×", 9.9, 5.34, 1.95, 0.45, 22, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, "lift vs portfolio", 9.9, 5.83, 1.95, 0.3, 10.5, GRAY, True, PP_ALIGN.CENTER)
    add_footer(slide, 11, "Source: Held-out test ranking. Campaign capacity should be tuned to action cost and contact-channel limits.")

    # 12. Retention playbook
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Business proposal: match risk signals to retention actions")
    headers = ["Risk pattern", "Recommended action", "Primary KPI"]
    widths = [3.35, 5.05, 2.75]
    x_positions = [0.78, 4.18, 9.28]
    for heading, x, width in zip(headers, x_positions, widths, strict=True):
        add_rect(slide, x, 1.42, width, 0.62, NAVY, radius=False)
        add_text(slide, heading, x + 0.16, 1.6, width - 0.32, 0.25, 12, WHITE, True)
    rows = [
        ("Ageing handset + high risk", "Upgrade offer or device-financing review", "Retention uplift"),
        ("Usage decline", "Plan-fit review, personalized bundle or usage benefit", "Usage recovery"),
        ("Service-quality friction", "Proactive care call and network/service resolution", "Complaint reduction"),
        ("High-value at-risk customer", "Priority outreach with value-capped incentive", "Net revenue retained"),
    ]
    colors = [PALE_BLUE, WHITE, PALE_BLUE, WHITE]
    for row_index, row in enumerate(rows):
        y = 2.05 + row_index * 0.94
        for value, x, width in zip(row, x_positions, widths, strict=True):
            add_rect(slide, x, y, width, 0.89, colors[row_index], radius=False, line="D9E2EC")
            add_text(slide, value, x + 0.16, y + 0.18, width - 0.32, 0.5, 11.5, INK, row_index == 0 and x == x_positions[0])
    add_rect(slide, 0.78, 6.08, 11.25, 0.54, "E6F6F2", line="B8E4D8")
    add_text(
        slide,
        "PoC design: randomly hold out 10% of each risk/action segment to measure incremental retention, not just response rate.",
        1.05,
        6.23,
        10.75,
        0.25,
        11,
        "176B58",
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 12, "Source: Proposed operating model derived from the leading predictive features.")

    # 13. Impact
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Illustrative economics support a controlled rollout", "Scenario per 100,000 scored customers; replace assumptions with Company A campaign costs")
    add_picture(slide, ASSET_DIR / "impact_scenarios.png", 0.65, 1.45, 7.25)
    add_rect(slide, 8.25, 1.55, 4.0, 4.9, LIGHT, line=PALE_BLUE)
    add_text(slide, "Base-case assumptions", 8.58, 1.87, 2.7, 0.4, 17, NAVY, True)
    assumptions = [
        "Target top 20% = 20,000 customers",
        f"Capture {analysis['capture']:.1%} of expected churners",
        f"Average target revenue = ${analysis['average_target_revenue']:.2f}/month",
        "20% of contacted churners retained",
        "$15 intervention cost per targeted customer",
        "12 months of retained revenue value",
    ]
    add_bullet_list(slide, assumptions, 8.52, 2.42, 3.2, 2.2, 11.5)
    add_line(slide, 8.58, 4.9, 11.85, 4.9, PALE_BLUE, 1)
    add_text(slide, "Net annual value", 8.58, 5.18, 1.8, 0.3, 11.5, GRAY, True)
    add_text(slide, f"${analysis['net_value']/1_000_000:.2f}M", 8.55, 5.55, 1.8, 0.52, 26, GREEN, True)
    add_text(slide, f"{analysis['roi']:.1f}× ROI", 10.45, 5.63, 1.35, 0.35, 17, BLUE, True, PP_ALIGN.RIGHT)
    add_footer(slide, 13, "Source: Model targeting results plus explicitly stated assumptions. This is a planning scenario, not measured campaign impact.")

    # 14. Roadmap
    slide = presentation.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "90-day PoC roadmap with decision gates")
    phases = [
        ("Weeks 1–2", "Validate", "Confirm data freshness, customer eligibility, costs and action catalogue.", BLUE),
        ("Weeks 3–5", "Integrate", "Automate weekly scoring and connect ranked lists to CRM workflows.", CYAN),
        ("Weeks 6–9", "Experiment", "Run segmented treatment/control campaigns and monitor operational quality.", ORANGE),
        ("Weeks 10–12", "Decide", "Measure incremental retention, net value, fairness and model drift.", GREEN),
    ]
    for index, (timing, heading, body, color) in enumerate(phases):
        x = 0.75 + index * 3.0
        add_rect(slide, x, 1.55, 2.7, 3.0, WHITE, line=PALE_BLUE)
        add_rect(slide, x, 1.55, 2.7, 0.18, color, radius=False)
        add_text(slide, timing, x + 0.22, 1.95, 1.3, 0.3, 10.5, color, True)
        add_text(slide, heading, x + 0.22, 2.4, 2.0, 0.4, 18, NAVY, True)
        add_text(slide, body, x + 0.22, 3.05, 2.18, 1.05, 11.5, GRAY)
    add_rect(slide, 0.75, 5.05, 11.7, 1.28, NAVY)
    add_text(slide, "Scale only if the PoC passes all four gates", 1.0, 5.36, 3.2, 0.35, 16, CYAN, True)
    gates = ["Incremental retention", "Positive net value", "Stable calibration", "No material segment harm"]
    for index, gate in enumerate(gates):
        x = 4.35 + index * 1.92
        add_rect(slide, x, 5.31, 1.65, 0.63, "173F5F")
        add_text(slide, gate, x + 0.08, 5.48, 1.48, 0.25, 9.5, WHITE, True, PP_ALIGN.CENTER)
    add_footer(slide, 14, "Source: Proposed implementation roadmap and model-risk controls.")

    # 15. Decision and references
    slide = presentation.slides.add_slide(blank)
    set_background(slide, NAVY)
    add_text(slide, "Decision requested", 0.72, 0.5, 5.0, 0.5, 26, WHITE, True)
    add_rect(slide, 0.72, 1.22, 5.45, 2.32, "173F5F")
    add_text(
        slide,
        "Approve a 90-day retention PoC targeting the highest-risk 20% of eligible customers, with randomized controls and value-capped interventions.",
        1.02,
        1.5,
        4.85,
        1.75,
        18,
        WHITE,
        True,
    )
    add_text(slide, "Expected PoC outputs", 0.85, 3.9, 2.5, 0.35, 14, CYAN, True)
    add_bullet_list(
        slide,
        [
            "Incremental churn reduction by segment and action",
            "Verified campaign ROI and scalable capacity",
            "Production monitoring and retraining requirements",
        ],
        0.82,
        4.38,
        5.2,
        1.45,
        13,
        WHITE,
    )
    add_line(slide, 6.65, 0.6, 6.65, 6.85, "315A73", 1.5)
    add_text(slide, "References", 7.05, 0.5, 4.5, 0.5, 24, WHITE, True)
    references = [
        "[1] Deloitte. 2026 Telecommunications Industry Outlook. https://www.deloitte.com/us/en/insights/industry/technology/technology-media-and-telecom-predictions/2026/telecommunications-industry-outlook.html",
        "[2] Company A Dataset Overview. ENG_Company.md. Supplied course material.",
        "[3] GCI World 2026 Final Assignment Tutorial. tutorial.ipynb. Supplied course material.",
        "[4] Company A analysis pipeline: train_models.py, train_improved_model.py, train_advanced_model.py.",
        "[5] Held-out metrics and figures: artifacts/advanced/metrics.json and feature_importance.csv.",
        "[6] Sample presentation used for narrative reference: Sample1.pptx.",
    ]
    y = 1.28
    for reference in references:
        add_text(slide, reference, 7.05, y, 5.35, 0.7, 10.2, "D9EAF4")
        y += 0.82
    add_text(
        slide,
        "All financial values on slide 13 are explicitly labeled scenario assumptions.",
        7.05,
        6.38,
        5.2,
        0.35,
        10,
        ORANGE,
        True,
    )
    add_text(slide, "15", 12.0, 7.13, 0.55, 0.2, 8, "AFC7D5", True, PP_ALIGN.RIGHT)

    presentation.save(OUTPUT_PATH)
    print(f"Created {OUTPUT_PATH}")


def main() -> None:
    analysis = load_analysis()
    generate_assets(analysis)
    build_deck(analysis)


if __name__ == "__main__":
    main()
